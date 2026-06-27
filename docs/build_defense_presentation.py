"""
Pulsefy diploma defense presentation — dark mode, 12 slides, ~5 minutes.

More descriptive text per slide than the first version, restrained academic
tone (no selling language), darker palette for sun-lit rooms.

Run from the worktree root:
    python3 docs/build_defense_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OUT = (
    "/Users/dumitruvlad/Documents/Facultate/Licenta/Pulsefy/"
    ".claude/worktrees/friendly-shockley-28309b/docs/"
    "Pulsefy_Defense_Presentation.pptx"
)

# Dark-mode palette
BG = RGBColor(0x0F, 0x17, 0x2A)        # slate-900 background
CARD = RGBColor(0x1E, 0x29, 0x3B)      # slate-800 card surface
INK = RGBColor(0xF8, 0xFA, 0xFC)       # slate-50 primary text
MUTED = RGBColor(0x94, 0xA3, 0xB8)     # slate-400 muted text
VIOLET = RGBColor(0xA7, 0x8B, 0xFA)    # violet-400 accent
ORANGE = RGBColor(0xFB, 0x92, 0x3C)    # orange-400 secondary accent (premium)
GREEN = RGBColor(0x4A, 0xDE, 0x80)     # green-400 (free tier accent)
BORDER = RGBColor(0x33, 0x41, 0x55)    # slate-700 subtle borders

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]


def add_slide():
    slide = prs.slides.add_slide(BLANK)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    return slide


def text(slide, x, y, w, h, content, *, size=18, bold=False, color=INK,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False,
         font_name="Calibri"):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = content
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return tb


def addp(tb, content, *, size=16, bold=False, color=INK,
         align=PP_ALIGN.LEFT, italic=False, font_name="Calibri"):
    tf = tb.text_frame
    p = tf.add_paragraph()
    p.alignment = align
    run = p.add_run()
    run.text = content
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return p


def accent_bar(slide, x, y, w, h, color=VIOLET):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


def card(slide, x, y, w, h, fill=CARD, border=BORDER, border_w=1.0):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border
    shp.line.width = Pt(border_w)
    return shp


def image_placeholder(slide, x, y, w, h, caption):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = CARD
    shp.line.color.rgb = BORDER
    shp.line.width = Pt(1)
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = caption
    run.font.size = Pt(13)
    run.font.italic = True
    run.font.color.rgb = MUTED


def slide_header(slide, title_text):
    accent_bar(slide, 1.0, 0.95, 0.08, 0.45, VIOLET)
    text(slide, 1.25, 0.85, 11.5, 0.65, title_text,
         size=28, bold=True, color=INK)


def notes(slide, content):
    slide.notes_slide.notes_text_frame.text = content


# ====================================================================
# SLIDE 1 — Title
# ====================================================================
s = add_slide()
accent_bar(s, 0, 3.4, 13.333, 0.04, VIOLET)
text(s, 1.0, 2.3, 11.3, 1.1, "Pulsefy",
     size=80, bold=True, color=INK)
text(s, 1.0, 3.6, 11.3, 0.7,
     "AI-Assisted Ad Creation Platform for Short-Form Social Video",
     size=22, color=MUTED)

text(s, 1.0, 5.6, 5.5, 0.4, "Vlad DUMITRU", size=15, bold=True, color=INK)
text(s, 1.0, 6.0, 5.5, 0.4, "Diploma project, June 2026",
     size=12, color=MUTED)

text(s, 7.8, 5.6, 4.6, 0.4, "Coordinator",
     size=11, color=MUTED, align=PP_ALIGN.RIGHT)
text(s, 7.8, 5.9, 4.6, 0.4, "Conf. dr. ing. Iuliana MARIN",
     size=14, bold=True, color=INK, align=PP_ALIGN.RIGHT)
text(s, 7.8, 6.3, 4.6, 0.4,
     "Faculty of Engineering in Foreign Languages, UPB",
     size=10, color=MUTED, align=PP_ALIGN.RIGHT)

notes(s,
    "Good morning. I'm Vlad Dumitru, presenting my diploma project, "
    "Pulsefy — an AI-assisted ad creation platform built under the "
    "guidance of Conf. dr. ing. Iuliana Marin at the Faculty of "
    "Engineering in Foreign Languages."
)


# ====================================================================
# SLIDE 2 — Context
# ====================================================================
s = add_slide()
slide_header(s, "Context: short-form social video advertising")

text(s, 1.0, 1.9, 11.3, 1.6,
     "TikTok, Instagram Reels, and YouTube Shorts have become the dominant "
     "surface for product advertising. The format demands frequent posting, "
     "vertical video, and audio-first content. The production cost gap "
     "between professional agency work and what an individual creator can "
     "produce alone has widened.",
     size=15, color=INK)

card(s, 1.0, 3.7, 11.3, 3.2)
text(s, 1.4, 3.9, 10.5, 0.5,
     "Three options currently available to a small advertiser",
     size=16, bold=True, color=VIOLET)

text(s, 1.4, 4.6, 3.4, 0.4, "Hire an agency",
     size=15, bold=True, color=INK)
text(s, 1.4, 5.05, 3.4, 1.6,
     "Roughly one to ten thousand euros per campaign. Out of reach for "
     "recurring use.",
     size=12, color=MUTED)

text(s, 5.0, 4.6, 3.4, 0.4, "Stitch four DIY tools",
     size=15, bold=True, color=INK)
text(s, 5.0, 5.05, 3.4, 1.6,
     "Canva, Suno, ElevenLabs, CapCut. Four subscriptions, four "
     "interfaces, hours of friction per ad.",
     size=12, color=MUTED)

text(s, 8.6, 4.6, 3.4, 0.4, "Skip advertising",
     size=15, bold=True, color=INK)
text(s, 8.6, 5.05, 3.4, 1.6,
     "Lose reach on platforms whose algorithms reward consistent "
     "posting cadence.",
     size=12, color=MUTED)

notes(s,
    "Short-form social video is now the dominant advertising surface for "
    "small businesses and creators. The current options are narrow: hire "
    "an agency at thousands of euros per campaign, stitch together four "
    "separate creative tools at the cost of hours per ad, or skip "
    "advertising and lose visibility. This is the gap Pulsefy targets."
)


# ====================================================================
# SLIDE 3 — What Pulsefy is
# ====================================================================
s = add_slide()
slide_header(s, "What Pulsefy is")

text(s, 1.0, 2.0, 11.3, 2.8,
     "Pulsefy is a web application that unifies four AI generation "
     "subsystems behind a single user account: scene image generation, "
     "background music generation, voiceover synthesis, and video "
     "assembly. The user enters a short product brief and receives a "
     "finished MP4 ready for upload to TikTok, Instagram Reels, or "
     "YouTube Shorts.",
     size=18, color=INK)

card(s, 1.0, 4.9, 11.3, 1.9)
text(s, 1.4, 5.1, 10.5, 0.5, "Scope of the platform",
     size=14, bold=True, color=VIOLET)
text(s, 1.4, 5.6, 10.5, 1.1,
     "A free tier with functional coverage across all four subsystems, "
     "a premium tier that substitutes higher-quality engines for the "
     "same workflows, and a social layer (Jamendo catalog browser and "
     "creator leaderboard) that connects users to one another and to "
     "Creative Commons music sources.",
     size=13, color=INK)

notes(s,
    "Pulsefy unifies the four generation steps that normally live in "
    "separate apps into a single web application. The user enters a "
    "brief, the system produces a finished MP4 in the requested aspect "
    "ratio, and the same account covers both free and premium "
    "workflows plus a thin social layer with the Jamendo catalog "
    "and a creator leaderboard."
)


# ====================================================================
# SLIDE 4 — Core features
# ====================================================================
s = add_slide()
slide_header(s, "Core features")

def feature_card(x, y, w, h, title, body):
    card(s, x, y, w, h)
    text(s, x + 0.25, y + 0.15, w - 0.5, 0.4, title,
         size=13, bold=True, color=VIOLET)
    text(s, x + 0.25, y + 0.6, w - 0.5, h - 0.7, body,
         size=11, color=INK)

# Two columns × three rows
feature_card(1.0, 1.85, 5.5, 1.45,
    "AI scene image generation",
    "Pollinations.ai integration with four named ad shot types and "
    "three social aspect ratios (9:16, 1:1, 16:9).")
feature_card(6.8, 1.85, 5.5, 1.45,
    "Two music engines",
    "Pulsefy AI custom LSTM for instant free-tier output, MusicGen "
    "for higher-quality premium output.")

feature_card(1.0, 3.45, 5.5, 1.45,
    "Multilingual voiceover",
    "gTTS-based voiceover synthesis. English on the free tier, "
    "dozens of languages on premium.")
feature_card(6.8, 3.45, 5.5, 1.45,
    "Video assembly and upload",
    "FFmpeg-based assembly with Ken Burns effect, plus a separate "
    "upload path that overlays generated audio onto user video.")

feature_card(1.0, 5.05, 5.5, 1.45,
    "Licensed music catalog",
    "Jamendo Creative Commons proxy with search, genre filtering, "
    "inline preview, and direct attachment.")
feature_card(6.8, 5.05, 5.5, 1.45,
    "Creator leaderboard",
    "Public ranking by generation count, surfacing active creators "
    "without requiring an explicit follow graph.")

notes(s,
    "Six core features. The four AI subsystems on the left side of "
    "the slide handle generation. The two on the right cover catalog "
    "browsing and the social-layer leaderboard. Together they cover "
    "the full lifecycle of a short-form ad from brief to publish-ready "
    "asset, with no need to leave the application."
)


# ====================================================================
# SLIDE 5 — Generation pipeline
# ====================================================================
s = add_slide()
slide_header(s, "Generation pipeline")

def pbox(x, y, w, h, label, sub, fill=CARD):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = VIOLET
    shp.line.width = Pt(1.25)
    tf = shp.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.size = Pt(14)
    r.font.bold = True
    r.font.color.rgb = INK
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = sub
    r2.font.size = Pt(10)
    r2.font.italic = True
    r2.font.color.rgb = MUTED

pbox(0.6, 3.2, 2.0, 1.5, "Product brief", "name, style, ratio")
pbox(3.4, 1.7, 3.3, 1.1, "Scene images", "Pollinations.ai")
pbox(3.4, 3.0, 3.3, 1.1, "Music", "MusicGen / Pulsefy AI")
pbox(3.4, 4.3, 3.3, 1.1, "Voiceover", "gTTS")
pbox(3.4, 5.6, 3.3, 1.1, "Assembly", "FFmpeg + Ken Burns")
pbox(8.0, 3.2, 2.0, 1.5, "MP4 ad", "9:16 / 1:1 / 16:9")
pbox(10.8, 3.2, 2.0, 1.5, "Publish", "TikTok / Reels / Shorts",
     fill=RGBColor(0x16, 0x2F, 0x1F))

text(s, 1.0, 6.95, 11.3, 0.45,
     "Single request from the brief, four parallel asset stages, "
     "assembly into a single MP4 at the chosen social-platform ratio.",
     size=11, color=MUTED, italic=True, align=PP_ALIGN.CENTER)

notes(s,
    "The pipeline takes a single product brief, fans out into four "
    "parallel generation stages, and converges on FFmpeg assembly. "
    "The result is a single MP4 in the user's chosen aspect ratio, "
    "ready for upload to the target social platform."
)


# ====================================================================
# SLIDE 6 — Original contribution: Pulsefy AI
# ====================================================================
s = add_slide()
slide_header(s, "Original contribution: Pulsefy AI")

text(s, 1.0, 1.85, 7.0, 0.5,
     "A custom LSTM music generator trained from scratch",
     size=16, color=VIOLET, italic=True)

text(s, 1.0, 2.7, 7.0, 4.0,
     "Pulsefy AI is the original technical contribution of the project. "
     "It is a token-level next-event LSTM model that generates "
     "background music for the free tier of the platform without "
     "depending on any third-party service.",
     size=14, color=INK)

addp(s.shapes[-1], "")
addp(s.shapes[-1],
    "Building this rather than relying solely on Meta's MusicGen "
    "required end-to-end implementation of three components: a MIDI "
    "preprocessing and tokenisation pipeline, a PyTorch training loop "
    "with Apple Metal acceleration, and an inference path integrated "
    "with the Node.js backend through a Python subprocess.",
    size=14, color=INK)

addp(s.shapes[-1], "")
addp(s.shapes[-1],
    "Inference completes in single-digit seconds on consumer "
    "hardware, which makes the free tier feel instant.",
    size=14, color=INK, italic=True)

image_placeholder(s, 8.4, 2.5, 4.5, 4.3,
    "[Insert: training screenshot from\nthesis Figure 43]")

notes(s,
    "The most substantive original contribution is Pulsefy AI, a "
    "custom LSTM music generator I trained from scratch. Building "
    "this myself rather than only plugging in Meta's MusicGen meant "
    "I had to implement the full pipeline — MIDI preprocessing and "
    "tokenisation, the PyTorch training loop on Apple Metal, and "
    "subprocess integration with the Node backend. Inference is fast "
    "enough that the free tier feels instant, which was the design goal."
)


# ====================================================================
# SLIDE 7 — Training methodology
# ====================================================================
s = add_slide()
slide_header(s, "Training methodology and results")

text(s, 1.0, 1.85, 11.3, 0.5,
     "Pulsefy AI was trained on a curated MIDI corpus over thirty epochs.",
     size=14, color=MUTED, italic=True)

def stat_col(x, label, value, sub):
    text(s, x, 2.7, 3.0, 0.4, label,
         size=12, color=MUTED, align=PP_ALIGN.CENTER)
    text(s, x, 3.15, 3.0, 1.0, value,
         size=38, bold=True, color=VIOLET, align=PP_ALIGN.CENTER)
    text(s, x, 4.5, 3.0, 0.5, sub,
         size=11, color=MUTED, align=PP_ALIGN.CENTER, italic=True)

stat_col(0.4,  "Input corpus", "1,034",   "MIDI files (folk + Christmas)")
stat_col(3.6,  "Tokens",       "573,692", "after preprocessing")
stat_col(6.7,  "Parameters",   "908,708", "2 LSTM layers, 256 hidden")
stat_col(9.8,  "Epochs",       "30",      "on Apple Metal (MPS)")

card(s, 1.0, 5.3, 11.3, 1.6)
text(s, 1.4, 5.45, 10.5, 0.5,
     "Convergence",
     size=13, bold=True, color=VIOLET)
text(s, 1.4, 5.9, 10.5, 1.0,
     "Validation loss converged from approximately 0.51 at epoch 1 to "
     "approximately 0.43 by epoch 3, indicating the model fit the "
     "dataset early and continued to improve over the full training "
     "run. The final checkpoint ships with the application.",
     size=12, color=INK)

notes(s,
    "The training corpus was 1,034 MIDI files, parsed into 573,692 "
    "tokens, used to train an LSTM with roughly 908 thousand parameters "
    "across two recurrent layers. Training ran for 30 epochs on Apple "
    "Metal hardware. Validation loss converged from 0.51 to 0.43 in "
    "the first three epochs, which indicates the architecture fits "
    "the dataset cleanly. The final checkpoint ships with the application."
)


# ====================================================================
# SLIDE 8 — Tier model
# ====================================================================
s = add_slide()
slide_header(s, "Tier model: functional free, quality-gated premium")

text(s, 1.0, 1.9, 11.3, 0.5,
     "The premium tier does not unlock new feature categories; it "
     "substitutes higher-quality engines for the same workflows.",
     size=14, color=MUTED, italic=True)

def tier_card(x, header, color, body_lines):
    card(s, x, 2.9, 5.5, 3.8, border=color, border_w=2.0)
    text(s, x + 0.3, 3.1, 5.0, 0.5, header,
         size=20, bold=True, color=color)
    y = 3.8
    for line in body_lines:
        text(s, x + 0.3, y, 5.0, 0.6, line,
             size=13, color=INK)
        y += 0.7

tier_card(0.6, "Free  ·  Pulsefy AI", GREEN, [
    "Custom LSTM (908k parameters)",
    "Music generation in 3–7 seconds",
    "English voiceover via gTTS",
    "Three scene images per ad",
])
tier_card(7.2, "Premium  ·  MusicGen", ORANGE, [
    "Meta audiocraft (balanced quality)",
    "Music generation in 5–7 minutes",
    "Voiceover in dozens of languages",
    "Up to four scene images per ad",
])

notes(s,
    "The tier model is deliberately quality-gated rather than feature-"
    "gated. Free users get the Pulsefy AI engine, instant but "
    "acoustically simpler, with English voiceover and three scene "
    "images. Premium users substitute Meta's MusicGen, slower but "
    "higher quality, and unlock multilingual voiceover and the fourth "
    "scene. The free tier is functionally complete — it is not a trial."
)


# ====================================================================
# SLIDE 9 — Sound Studio screenshot
# ====================================================================
s = add_slide()
slide_header(s, "The Sound Studio")

text(s, 1.0, 1.85, 11.3, 0.9,
     "Four tabs correspond to the four AI subsystems. The user moves "
     "between them at their own pace and the underlying session state "
     "persists across navigations.",
     size=13, color=MUTED, italic=True)

image_placeholder(s, 1.0, 2.95, 11.3, 4.0,
    "[Insert: Sound Studio screenshot — Music tab with the\n"
    "Pulsefy AI vs MusicGen selector visible]")

notes(s,
    "This is the Sound Studio, the central authoring surface. Four "
    "tabs — Music, Voiceover, Video, and Upload — correspond to the "
    "four AI subsystems. The session state persists as the user moves "
    "between tabs, so they can generate the music and the voiceover "
    "separately and combine them in the Video tab without losing "
    "their work."
)


# ====================================================================
# SLIDE 10 — Performance benchmark
# ====================================================================
s = add_slide()
slide_header(s, "Performance benchmark")

text(s, 1.0, 1.85, 11.3, 0.5,
     "Three Instagram Reel scenarios, measured on Mac M2 Pro.",
     size=14, color=MUTED, italic=True)

def benchrow(y, label, value, sub, color):
    card(s, 1.0, y, 11.3, 0.85)
    text(s, 1.3, y + 0.05, 4.5, 0.75, label,
         size=15, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    text(s, 5.8, y + 0.05, 3.5, 0.75, value,
         size=22, bold=True, color=color,
         anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    text(s, 9.5, y + 0.05, 2.7, 0.75, sub,
         size=11, color=MUTED, italic=True,
         anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)

benchrow(2.85, "Manual workflow (Canva + Suno + ElevenLabs + CapCut)",
         "~8 hours", "baseline", MUTED)
benchrow(3.95, "Pulsefy free tier (Pulsefy AI music)",
         "~3 min 11 s", "instant feedback loop", GREEN)
benchrow(5.05, "Pulsefy premium tier (MusicGen music)",
         "~9 min 12 s", "MusicGen dominates the time", ORANGE)

text(s, 1.0, 6.3, 11.3, 0.6,
     "Quality assessment was out of scope for the thesis benchmark. "
     "A larger listening study with objective audio metrics is identified "
     "as future work.",
     size=11, color=MUTED, italic=True, align=PP_ALIGN.LEFT)

notes(s,
    "I benchmarked the platform on a Mac M2 Pro across three Instagram "
    "Reel scenarios. The free-tier pipeline completes a full ad in "
    "around three minutes; the premium tier takes about nine minutes, "
    "with MusicGen accounting for most of that time. Quality assessment "
    "of the generated content is out of scope for this benchmark — "
    "the goal here was to verify the latency claim from Chapter 1."
)


# ====================================================================
# SLIDE 11 — Limitations + Future work
# ====================================================================
s = add_slide()
slide_header(s, "Limitations and future work")

# Limitations column
text(s, 1.0, 1.85, 5.5, 0.5, "Current limitations",
     size=16, bold=True, color=ORANGE)
def lim(y, head, body):
    text(s, 1.0, y, 5.5, 0.35, head, size=12, bold=True, color=INK)
    text(s, 1.0, y + 0.35, 5.5, 0.7, body, size=11, color=MUTED)
lim(2.5, "Benchmark scope",
    "N=3, single hardware, latency-only (no quality assessment).")
lim(3.75, "MusicGen latency",
    "5–7 minutes per request on CPU is the slowest stage of the pipeline.")
lim(5.0, "Desktop-only UI",
    "The four-tab Sound Studio expects a desktop-class viewport.")
lim(6.25, "Simulated subscriptions",
    "Premium upgrade is self-service without a real payment processor.")

# Future work column
text(s, 7.0, 1.85, 5.5, 0.5, "Future work",
     size=16, bold=True, color=VIOLET)
def fw(y, head, body):
    text(s, 7.0, y, 5.5, 0.35, head, size=12, bold=True, color=INK)
    text(s, 7.0, y + 0.35, 5.5, 0.7, body, size=11, color=MUTED)
fw(2.5, "GPU-backed inference",
    "Worker pool on GPU hardware to reduce MusicGen latency to seconds.")
fw(3.75, "Larger validation study",
    "N≥20 listeners, MUSHRA scoring, Fréchet Audio Distance for quality.")
fw(5.0, "Mobile-responsive UI",
    "Meet creators on the device they shoot their content on.")
fw(6.25, "Production payment + multi-language UI",
    "Stripe integration; Romanian and other regional language support.")

notes(s,
    "Four current limitations: the benchmark is small, MusicGen latency "
    "dominates the premium pipeline, the UI is desktop-only, and "
    "premium upgrades are simulated. Four corresponding future-work "
    "items address each of these: GPU-backed inference for MusicGen, "
    "a larger listening study with objective quality metrics, a "
    "mobile-responsive layout, and integration with a real payment "
    "processor along with multi-language UI support."
)


# ====================================================================
# SLIDE 12 — Closing
# ====================================================================
s = add_slide()
accent_bar(s, 0, 3.4, 13.333, 0.04, VIOLET)
text(s, 0.5, 2.7, 12.3, 1.0, "Thank you for your attention.",
     size=44, bold=True, color=INK, align=PP_ALIGN.CENTER)
text(s, 0.5, 4.0, 12.3, 0.6, "I welcome your questions.",
     size=22, color=VIOLET, align=PP_ALIGN.CENTER, italic=True)
text(s, 0.5, 6.7, 12.3, 0.4,
     "Vlad DUMITRU  ·  Pulsefy  ·  UPB FILS 2026",
     size=11, color=MUTED, align=PP_ALIGN.CENTER)

notes(s,
    "Thank you for your attention. I welcome your questions."
)


prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides)}")
