"""
Surgical edits to the existing Pulsefy_Defense_Presentation.pptx.

Opens Vlad's current .pptx (preserving any manual edits / inserted screenshots)
and applies five changes:
  1. Icons (small violet star) on slides 3, 4, 6, 7, 10
  2. Connector arrows on slide 5 (pipeline)
  3. Replace slide 7 (training) with a simpler 3-step flow
  4. Insert a new "Who Pulsefy is for" slide between slides 3 and 4
  5. Update slide 11 (performance benchmark, was 10) with the 8-hour breakdown

This script does NOT regenerate from scratch — it opens the existing file
and modifies in place. Run from the worktree root:

    python3 docs/edit_defense_presentation.py
"""

from copy import deepcopy
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
from lxml import etree

PPTX = (
    "/Users/dumitruvlad/Documents/Facultate/Licenta/Pulsefy/"
    ".claude/worktrees/friendly-shockley-28309b/docs/"
    "Pulsefy_Defense_Presentation.pptx"
)

# Palette (same as v2 dark theme)
BG = RGBColor(0x0F, 0x17, 0x2A)
CARD = RGBColor(0x1E, 0x29, 0x3B)
INK = RGBColor(0xF8, 0xFA, 0xFC)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
VIOLET = RGBColor(0xA7, 0x8B, 0xFA)
GREEN = RGBColor(0x4A, 0xDE, 0x80)
ORANGE = RGBColor(0xFB, 0x92, 0x3C)
BORDER = RGBColor(0x33, 0x41, 0x55)


prs = Presentation(PPTX)


def text(slide, x, y, w, h, content, *, size=14, bold=False, color=INK,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False,
         font_name="Calibri"):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = content
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return tb


def addp(tb, content, *, size=12, bold=False, color=INK, align=PP_ALIGN.LEFT,
         italic=False, font_name="Calibri"):
    p = tb.text_frame.add_paragraph()
    p.alignment = align
    run = p.add_run()
    run.text = content
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return p


def rounded_card(slide, x, y, w, h, fill=CARD, border=BORDER, border_w=1.0):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border
    shp.line.width = Pt(border_w)
    return shp


def add_icon_star(slide, x=12.4, y=0.7, size=0.7, color=VIOLET, glyph="★"):
    """Place a small violet symbol top-right as an icon accent."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(size), Inches(size))
    tf = tb.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = glyph
    r.font.size = Pt(28)
    r.font.color.rgb = color
    return tb


def add_arrow(slide, x1, y1, x2, y2, color=VIOLET, weight=1.5):
    """Draw a straight line connector with a triangle arrowhead at the destination."""
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                       Inches(x1), Inches(y1),
                                       Inches(x2), Inches(y2))
    ln = conn.line
    ln.color.rgb = color
    ln.width = Pt(weight)
    # XML hack to add a triangle arrowhead at the END of the line
    ln_xml = conn.line._get_or_add_ln()
    tail = etree.SubElement(
        ln_xml, qn('a:tailEnd'),
        {'type': 'triangle', 'w': 'med', 'len': 'med'}
    )
    return conn


def clear_slide_shapes(slide):
    """Remove every shape from a slide (preserving slide background)."""
    spTree = slide.shapes._spTree
    for shp in list(slide.shapes):
        spTree.remove(shp._element)


def insert_slide_at(prs, slide, position):
    """Move a slide (already added to end) to position `position` (0-indexed)."""
    xml_slides = prs.slides._sldIdLst
    children = list(xml_slides)
    # The newly added slide is the last child
    moved = children[-1]
    xml_slides.remove(moved)
    xml_slides.insert(position, moved)


# ─── 1. Icons on slides 3, 4, 6, 7, 10 (1-indexed) ───────────────────────
# After we insert the new "Who Pulsefy is for" slide between 3 and 4 (step 4),
# the numbering shifts. To avoid juggling, we apply icons FIRST against
# the current (pre-insertion) numbering: slides 3, 4, 6, 7, 10.
# 0-indexed: 2, 3, 5, 6, 9
icon_targets = [2, 3, 5, 6, 9]
for idx in icon_targets:
    if idx < len(prs.slides):
        add_icon_star(prs.slides[idx])
print(f"Icons added on slides {[i+1 for i in icon_targets]}")


# ─── 2. Arrows on slide 5 (pipeline) ──────────────────────────────────────
# Coordinates from the original v2 layout:
#   Brief box:     (0.6, 3.2, 2.0, 1.5)  → right edge x=2.6, mid y=3.95
#   Scene images:  (3.4, 1.7, 3.3, 1.1)  → left x=3.4, mid y=2.25
#   Music:         (3.4, 3.0, 3.3, 1.1)  → left x=3.4, mid y=3.55
#   Voiceover:     (3.4, 4.3, 3.3, 1.1)  → left x=3.4, mid y=4.85
#   Assembly:      (3.4, 5.6, 3.3, 1.1)  → left x=3.4, mid y=6.15
#   MP4 ad:        (8.0, 3.2, 2.0, 1.5)  → left x=8.0, mid y=3.95
#   Publish:       (10.8, 3.2, 2.0, 1.5) → left x=10.8, mid y=3.95
slide5 = prs.slides[4]
# Fan-out from brief
add_arrow(slide5, 2.6, 3.95, 3.35, 2.25)  # to scene images
add_arrow(slide5, 2.6, 3.95, 3.35, 3.55)  # to music
add_arrow(slide5, 2.6, 3.95, 3.35, 4.85)  # to voiceover
add_arrow(slide5, 2.6, 3.95, 3.35, 6.15)  # to assembly
# Fan-in to MP4 (right edge of generators is x=6.7)
add_arrow(slide5, 6.7, 2.25, 8.0, 3.95)
add_arrow(slide5, 6.7, 3.55, 8.0, 3.95)
add_arrow(slide5, 6.7, 4.85, 8.0, 3.95)
add_arrow(slide5, 6.7, 6.15, 8.0, 3.95)
# MP4 to Publish (right edge of MP4 is x=10.0)
add_arrow(slide5, 10.0, 3.95, 10.8, 3.95)
print("9 arrows added on slide 5")


# ─── 3. Replace slide 7 (training) with a simpler 3-step flow ─────────────
slide7 = prs.slides[6]
# Reapply dark background (in case it was reset when shapes cleared)
slide7.background.fill.solid()
slide7.background.fill.fore_color.rgb = BG
clear_slide_shapes(slide7)

# Re-add: accent bar + title
acc = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(1.0), Inches(0.95), Inches(0.08), Inches(0.45))
acc.fill.solid()
acc.fill.fore_color.rgb = VIOLET
acc.line.fill.background()

text(slide7, 1.25, 0.85, 11.5, 0.65,
     "How Pulsefy AI was trained",
     size=28, bold=True, color=INK)

# 3-step flow boxes
def step_box(x, y, w, h, label, sub):
    shp = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = CARD
    shp.line.color.rgb = VIOLET
    shp.line.width = Pt(1.5)
    tf = shp.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = INK
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = sub
    r2.font.size = Pt(12)
    r2.font.italic = True
    r2.font.color.rgb = MUTED

# Three big boxes across the middle
step_box(0.7, 3.0, 3.7, 1.8,
         "1,034 MIDI files", "Nottingham folk + Christmas tunes")
step_box(4.8, 3.0, 3.7, 1.8,
         "Train for 30 epochs", "Apple Metal acceleration")
step_box(8.9, 3.0, 3.7, 1.8,
         "Pulsefy AI model", "saved as a .pt checkpoint")

# Arrows between boxes
add_arrow(slide7, 4.4, 3.9, 4.8, 3.9)  # step 1 → 2
add_arrow(slide7, 8.5, 3.9, 8.9, 3.9)  # step 2 → 3

# Bottom caption
text(slide7, 1.0, 5.5, 11.3, 1.0,
     "The model learned to predict the next musical event in a sequence — "
     "the same way a language model predicts the next word.",
     size=14, color=INK, italic=True, align=PP_ALIGN.CENTER)

# Add icon top-right (we already added the star earlier, but it got cleared
# when we cleared shapes; re-add)
add_icon_star(slide7)

# Speaker notes update
slide7.notes_slide.notes_text_frame.text = (
    "I trained Pulsefy AI from scratch on a curated dataset of 1,034 MIDI "
    "files — Nottingham folk waltzes and Christmas tunes. Training ran for "
    "thirty epochs on my Mac, accelerated by Apple Metal. The end result is "
    "a saved model file that the application ships with. Conceptually, the "
    "model learned the same way a small language model learns text: by "
    "predicting one token at a time, getting feedback on whether the "
    "prediction was right, and adjusting its internal weights over thousands "
    "of iterations."
)
print("Slide 7 rebuilt as simpler 3-step training flow")


# ─── 4. Insert new "Who Pulsefy is for" slide between slides 3 and 4 ─────
# Use the same blank layout (index 6).
BLANK = prs.slide_layouts[6]
new_slide = prs.slides.add_slide(BLANK)
new_slide.background.fill.solid()
new_slide.background.fill.fore_color.rgb = BG

# Title bar
acc2 = new_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(1.0), Inches(0.95),
                                  Inches(0.08), Inches(0.45))
acc2.fill.solid()
acc2.fill.fore_color.rgb = VIOLET
acc2.line.fill.background()
text(new_slide, 1.25, 0.85, 11.5, 0.65,
     "Who Pulsefy is for",
     size=28, bold=True, color=INK)

text(new_slide, 1.0, 1.85, 11.3, 0.5,
     "Three archetypes who produce ads themselves rather than commissioning "
     "them.",
     size=14, color=MUTED, italic=True)

# Three persona cards (horizontal)
def persona(x, header, body):
    rounded_card(new_slide, x, 2.8, 3.85, 3.7)
    text(new_slide, x + 0.3, 3.0, 3.4, 0.5, header,
         size=16, bold=True, color=VIOLET)
    text(new_slide, x + 0.3, 3.7, 3.4, 2.7, body,
         size=12, color=INK)

persona(0.7,
    "Solo creator",
    "A side-hustle TikTok or Reels creator producing promotional content "
    "for personal projects or sponsorships. Needs frequent output on a "
    "near-zero budget."
)
persona(4.75,
    "Small business owner",
    "A retail, food, or services owner who advertises on Reels and TikTok "
    "but cannot afford an agency or hire a full-time designer."
)
persona(8.8,
    "In-house marketer",
    "An early-stage startup employee who handles ads as one of many "
    "responsibilities. Needs to ship variants fast without learning four "
    "creative tools."
)

text(new_slide, 1.0, 6.7, 11.3, 0.4,
     "None can justify an agency budget; all need to ship more ads, more often.",
     size=12, color=MUTED, italic=True, align=PP_ALIGN.CENTER)

def set_notes(slide, text):
    ns = slide.notes_slide
    tf = ns.notes_text_frame
    if tf is None:
        # The slide layout had no notes placeholder; one exists once we
        # access notes_slide, but notes_text_frame may be None on a brand-new
        # slide. Fallback: find/create a notes_placeholder via the layout.
        for ph in ns.placeholders:
            if ph.has_text_frame:
                ph.text_frame.text = text
                return
        # Last-resort: skip silently
        return
    tf.text = text

set_notes(new_slide,
    "Pulsefy is built for users who produce ads themselves rather than "
    "commissioning them. Three concrete archetypes: a solo creator running "
    "their own social account, a small business owner handling their own "
    "marketing, and a startup employee whose job description includes ads "
    "alongside everything else. None of them can justify an agency budget "
    "or a multi-tool creative suite, and all of them need to ship more ads "
    "more often than the manual workflow allows."
)

# Move it to position 3 (0-indexed) so it becomes slide 4
insert_slide_at(prs, new_slide, 3)
print("New 'Who Pulsefy is for' slide inserted at position 4")


# ─── 5. Performance slide — add the 8-hour breakdown ──────────────────────
# After the insertion above, slide indices shifted by +1 for slides after #4.
# The performance benchmark was slide 10 → now slide 11 (0-indexed 10).
perf_slide = prs.slides[10]

# Find the existing caveat text box ("Quality assessment was out of scope...")
# and replace its text with the new combined version.
caveat_found = False
for shape in perf_slide.shapes:
    if not shape.has_text_frame:
        continue
    full_text = "\n".join(p.text for p in shape.text_frame.paragraphs)
    if "Quality assessment" in full_text or "quality assessment" in full_text:
        tf = shape.text_frame
        # Clear all paragraphs in the shape
        # python-pptx doesn't have a clean way, so we clear runs paragraph by paragraph
        for p in list(tf.paragraphs):
            for r in list(p.runs):
                r._r.getparent().remove(r._r)
        # First paragraph: breakdown header
        first = tf.paragraphs[0]
        r = first.add_run()
        r.text = "How the 8-hour baseline was calculated:"
        r.font.size = Pt(11)
        r.font.bold = True
        r.font.color.rgb = MUTED
        # Second paragraph: breakdown body
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = (
            "2 h product visuals in Canva + 2 h music in Suno or stock "
            "licensing + 1 h voiceover in ElevenLabs + 3 h editing and "
            "assembly in CapCut (breakdown from Chapter 1.2 of the thesis)."
        )
        r2.font.size = Pt(11)
        r2.font.color.rgb = MUTED
        # Third paragraph: quality caveat
        p3 = tf.add_paragraph()
        r3 = p3.add_run()
        r3.text = (
            "Quality assessment of generated content was out of scope; "
            "identified as future work."
        )
        r3.font.size = Pt(10)
        r3.font.italic = True
        r3.font.color.rgb = MUTED
        caveat_found = True
        break

if caveat_found:
    print("Performance slide updated with 8-hour breakdown")
else:
    print("WARN: caveat text box not found on performance slide; "
          "adding breakdown as new text box instead")
    text(perf_slide, 1.0, 6.3, 11.3, 1.0,
         "How the 8-hour baseline was calculated:",
         size=11, bold=True, color=MUTED, italic=True)
    text(perf_slide, 1.0, 6.55, 11.3, 0.4,
         "2 h product visuals in Canva + 2 h music in Suno or stock "
         "licensing + 1 h voiceover in ElevenLabs + 3 h editing and "
         "assembly in CapCut (breakdown from Chapter 1.2 of the thesis).",
         size=11, color=MUTED)


prs.save(PPTX)
print(f"\nSaved: {PPTX}")
print(f"Total slides: {len(prs.slides)}")
